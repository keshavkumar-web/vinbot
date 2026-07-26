"""Generate UHBVN_Assistant_Demo.pptx from the demo outline.

Run:  python make_pptx.py
Needs: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- Theme ----------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x2A, 0x4A)      # titles / bars
BLUE = RGBColor(0x1F, 0x6F, 0xB2)      # accents
GREEN = RGBColor(0x1E, 0x88, 0x5E)     # positive
RED = RGBColor(0xC0, 0x39, 0x2B)       # demo / alert
DARK = RGBColor(0x22, 0x2A, 0x33)      # body text
GREY = RGBColor(0x5A, 0x63, 0x6E)      # subtle
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)     # panel bg

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def set_run(run, text, size, color, bold=False, italic=False, font="Segoe UI"):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def header(slide, kicker, title):
    """Top navy bar with kicker + title."""
    rect(slide, 0, 0, SW, Inches(1.25), NAVY)
    rect(slide, 0, Inches(1.25), SW, Inches(0.06), BLUE)
    tf = textbox(slide, Inches(0.6), Inches(0.18), Inches(12.1), Inches(1.0))
    p = tf.paragraphs[0]
    set_run(p.add_run(), kicker.upper(), 12, RGBColor(0x9F, 0xC2, 0xE0), bold=True)
    p2 = tf.add_paragraph()
    set_run(p2.add_run(), title, 30, WHITE, bold=True)


def bullets(slide, items, x=Inches(0.7), y=Inches(1.7), w=Inches(11.9),
            h=Inches(5.3), size=20, gap=10):
    tf = textbox(slide, x, y, w, h)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        level = 0
        text = item
        color = DARK
        bold = False
        if isinstance(item, tuple):
            text, opts = item
            level = opts.get("level", 0)
            color = opts.get("color", DARK)
            bold = opts.get("bold", False)
            size_i = opts.get("size", size)
        else:
            size_i = size
        p.level = level
        bullet = "•  " if level == 0 else "–  "
        set_run(p.add_run(), bullet + text, size_i, color, bold=bold)
    return tf


def footer(slide, n):
    tf = textbox(slide, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4))
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Vinbot  ·  Demo", 10, GREY)
    p.alignment = PP_ALIGN.LEFT
    tf2 = textbox(slide, Inches(12.0), Inches(7.0), Inches(0.9), Inches(0.4))
    p2 = tf2.paragraphs[0]
    set_run(p2.add_run(), str(n), 10, GREY)
    p2.alignment = PP_ALIGN.RIGHT


# ============================================================== SLIDE 1 (title)
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.9), SW, Inches(0.08), BLUE)
tf = textbox(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.5))
p = tf.paragraphs[0]
set_run(p.add_run(), "Vinbot", 54, WHITE, bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(),
        "An AI chatbot that answers from UHBVN's sales circulars, "
        "compendiums & press notes", 22, RGBColor(0xBF, 0xD6, 0xEC))
tf3 = textbox(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.2))
p3 = tf3.paragraphs[0]
set_run(p3.add_run(), "Presented by: __________________", 16,
        RGBColor(0x9F, 0xC2, 0xE0))
p4 = tf3.add_paragraph()
set_run(p4.add_run(), "Date: __________________", 16,
        RGBColor(0x9F, 0xC2, 0xE0))

# ============================================================== SLIDE 2
s = add_slide()
header(s, "Why this matters", "The Problem")
bullets(s, [
    "UHBVN policy lives across hundreds of pages: sales circulars, "
    "compendiums of instructions, and press notes (2018 → 2026).",
    "Staff & consumers struggle to find the right clause quickly.",
    "Searching PDFs is slow; interpreting them needs expertise.",
    "Wrong or outdated information leads to disputes and rework.",
    ('Example: "What is the subsidized tariff for a registered Gaushala?" '
     "means knowing the circular, then reading it.",
     {"color": GREY, "bold": True, "size": 18}),
])
footer(s, 2)

# ============================================================== SLIDE 3
s = add_slide()
header(s, "What we built", "The Solution")
bullets(s, [
    "A chat interface: ask in plain language, get a grounded answer.",
    ("Answers come straight from official UHBVN documents.", {"level": 1}),
    ("The source circular / clause is cited in every answer.", {"level": 1}),
    ('Says "I don\'t have information about that" when it\'s not in the docs '
     "— no guessing.", {"level": 1}),
    ("", {}),
    ("Available 24/7   ·   Consistent answers   ·   Always cites the source",
     {"color": GREEN, "bold": True, "size": 22}),
])
footer(s, 3)

# ============================================================== SLIDE 4
s = add_slide()
header(s, "How it works", "Retrieval-Augmented Generation (in plain English)")
panel = rect(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(2.5), LIGHT)
tf = textbox(s, Inches(1.0), Inches(1.9), Inches(11.3), Inches(2.2))
flow = [
    "UHBVN PDFs  →  Extracted text  →  Split into chunks  →  Embeddings (knowledge DB)",
    "User question  →  Embedding  →  Find most similar chunks",
    "AI writes an answer using ONLY those chunks  →  User",
]
for i, line in enumerate(flow):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    set_run(p.add_run(), line, 16, NAVY, bold=True, font="Consolas")
bullets(s, [
    ("Key point: the AI is constrained to the knowledge base — it cannot "
     "invent figures or regulations.", {"color": RED, "bold": True, "size": 19}),
    "Unlike a generic chatbot, it only speaks from UHBVN's own circulars — "
    "and tells you which one.",
], y=Inches(4.5), h=Inches(2.2))
footer(s, 4)

# ============================================================== SLIDE 5
s = add_slide()
header(s, "Knowledge base", "What It Knows Today")
bullets(s, [
    ("55 official UHBVN documents already loaded, including:",
     {"bold": True, "size": 21}),
    ("Sales Circulars (U-series & I-series), 2018–2026", {"level": 1}),
    ("Compendium of Instructions (2022, 2023, 2025)", {"level": 1}),
    ("Press Notes & CGRF notes", {"level": 1}),
    ("RAPDRP / GIS module references", {"level": 1}),
    ("", {}),
    ("Easily extended — drop in a new PDF and re-run one command (see "
     "“Keeping It Updated”).", {"color": GREEN, "bold": True, "size": 19}),
])
footer(s, 5)

# ============================================================== SLIDE 6 (DEMO)
s = add_slide()
rect(s, 0, 0, SW, Inches(1.25), RED)
rect(s, 0, Inches(1.25), SW, Inches(0.06), RGBColor(0x8E, 0x24, 0x1B))
tf = textbox(s, Inches(0.6), Inches(0.18), Inches(12.1), Inches(1.0))
p = tf.paragraphs[0]
set_run(p.add_run(), "LIVE DEMO", 12, RGBColor(0xF5, 0xC6, 0xC0), bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(), "See It In Action", 30, WHITE, bold=True)
bullets(s, [
    ('1.  Simple lookup — "What is the subsidized electricity tariff for '
     'registered Gaushalas?"', {"bold": True, "size": 19}),
    ("Show the answer + the cited circular (SC U-01/2026).",
     {"level": 1, "color": GREY}),
    ('2.  Follow-up (memory) — "When did that take effect?"',
     {"bold": True, "size": 19}),
    ("Shows it remembers context within the conversation.",
     {"level": 1, "color": GREY}),
    ('3.  Guardrail — "Who won the cricket match yesterday?"',
     {"bold": True, "size": 19}),
    ('Politely refuses: "I don\'t have information about that…"',
     {"level": 1, "color": GREY}),
    ("4.  Citation + streaming — point out the [Source: …] reference and the "
     "live token-by-token typing.", {"bold": True, "size": 19}),
    ("Tip: pre-type these questions. If the API/network drops, use the "
     "backup screenshots.", {"color": RED, "bold": True, "size": 17}),
], y=Inches(1.6), gap=6)
footer(s, 6)

# ============================================================== SLIDE 7
s = add_slide()
header(s, "Backup", "Demo Screenshots (fallback)")
bullets(s, [
    ("Paste 2–3 screenshots here in case the live demo can't run:",
     {"bold": True, "size": 20}),
    ("A grounded answer with a citation", {"level": 1}),
    ('The "no information" guardrail response', {"level": 1}),
    ("The chat UI on screen", {"level": 1}),
])
rect(s, Inches(0.7), Inches(3.4), Inches(11.9), Inches(3.0),
     RGBColor(0xE6, 0xEA, 0xEF))
tf = textbox(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(0.8))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "[ screenshot placeholder ]", 18, GREY, italic=True)
footer(s, 7)

# ============================================================== SLIDE 8
s = add_slide()
header(s, "For technical reviewers", "Under the Hood")
rows = [
    ("Frontend", "Vue 3 + Vite + Tailwind"),
    ("Backend", "FastAPI (Python), streaming responses (SSE)"),
    ("AI models", "OpenAI gpt-4o-mini (chat) + text-embedding-3-small (search)"),
    ("Retrieval", "Cosine similarity over embedded document chunks"),
    ("Sessions", "Per-user chat memory (last 20 messages)"),
]
top = Inches(1.8)
rowh = Inches(0.78)
from pptx.enum.shapes import MSO_SHAPE
for i, (k, v) in enumerate(rows):
    y = top + rowh * i
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.7), y, Inches(11.9), rowh, bg)
    rect(s, Inches(0.7), y, Inches(3.0), rowh, NAVY)
    t1 = textbox(s, Inches(0.85), y + Inches(0.16), Inches(2.8), rowh)
    set_run(t1.paragraphs[0].add_run(), k, 16, WHITE, bold=True)
    t2 = textbox(s, Inches(3.9), y + Inches(0.16), Inches(8.5), rowh)
    set_run(t2.paragraphs[0].add_run(), v, 16, DARK)
tf = textbox(s, Inches(0.7), top + rowh * 5 + Inches(0.1), Inches(11.9), Inches(0.8))
set_run(tf.paragraphs[0].add_run(),
        "Clean separation of config, retrieval, chat & sessions. Runs as one "
        "server in production (API + UI on one origin).", 15, GREY, italic=True)
footer(s, 8)

# ============================================================== SLIDE 9
s = add_slide()
header(s, "Maintenance", "Keeping It Up to Date")
panel = rect(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.9),
             RGBColor(0x1E, 0x29, 0x33))
tf = textbox(s, Inches(1.0), Inches(1.85), Inches(11.3), Inches(1.6))
code = [
    "1. Drop the new PDF into the UHBVN/ folder",
    "2. Run:  python ingest.py --pdf-dir ../UHBVN --append",
    "3. Restart the backend",
]
for i, line in enumerate(code):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    set_run(p.add_run(), line, 16, RGBColor(0x8E, 0xE6, 0xB0), font="Consolas")
bullets(s, [
    "Extracts text → chunks → embeds → merges into the knowledge base.",
    "Only processes new files (fast, low cost).",
    "No code changes and no AI retraining needed.",
    ("Message for management: policy changes? The bot is current the same "
     "day — no developer needed.", {"color": GREEN, "bold": True, "size": 19}),
], y=Inches(3.9), h=Inches(2.8))
footer(s, 9)

# ============================================================== SLIDE 10
s = add_slide()
header(s, "Trust", "Trust & Safety")
bullets(s, [
    "Grounded answers only — refuses anything outside UHBVN documents.",
    "Always cites the source circular / clause.",
    "No hallucinated figures — forbidden from inventing charges or regulations.",
    "Stays on-topic — declines politics, general knowledge, etc.",
    "API keys kept in environment config, never in code (recently hardened).",
    "No consumer personal data stored — documents are public circulars.",
])
footer(s, 10)

# ============================================================== SLIDE 11
s = add_slide()
header(s, "Value", "Benefits & Impact")
# Two columns
rect(s, Inches(0.7), Inches(1.8), Inches(5.8), Inches(4.6), LIGHT)
rect(s, Inches(6.75), Inches(1.8), Inches(5.85), Inches(4.6),
     RGBColor(0xEC, 0xF4, 0xEC))
rect(s, Inches(0.7), Inches(1.8), Inches(5.8), Inches(0.7), NAVY)
rect(s, Inches(6.75), Inches(1.8), Inches(5.85), Inches(0.7), GREEN)
set_run(textbox(s, Inches(0.95), Inches(1.92), Inches(5.4), Inches(0.5))
        .paragraphs[0].add_run(), "For Consumers", 20, WHITE, bold=True)
set_run(textbox(s, Inches(7.0), Inches(1.92), Inches(5.4), Inches(0.5))
        .paragraphs[0].add_run(), "For UHBVN Staff", 20, WHITE, bold=True)
cons = ["Instant, plain-language answers", "24/7 availability",
        "Less dependence on office visits", "Reduced disputes"]
staff = ["Faster clause lookup", "Consistent, citable responses",
         "Less time searching PDFs", "Reduced training overhead"]
t1 = textbox(s, Inches(1.0), Inches(2.7), Inches(5.2), Inches(3.5))
for i, c in enumerate(cons):
    p = t1.paragraphs[0] if i == 0 else t1.add_paragraph()
    p.space_after = Pt(14)
    set_run(p.add_run(), "•  " + c, 18, DARK)
t2 = textbox(s, Inches(7.05), Inches(2.7), Inches(5.3), Inches(3.5))
for i, c in enumerate(staff):
    p = t2.paragraphs[0] if i == 0 else t2.add_paragraph()
    p.space_after = Pt(14)
    set_run(p.add_run(), "•  " + c, 18, DARK)
footer(s, 11)

# ============================================================== SLIDE 12
s = add_slide()
header(s, "Future", "Roadmap / What's Next")
bullets(s, [
    "Scale: move to a database / vector DB for many concurrent users.",
    "Coverage: auto-ingest new circulars from the UHBVN website.",
    "Access: deploy on UHBVN portal / WhatsApp / IVR.",
    "Languages: Hindi + Haryanvi support.",
    "Analytics: track top questions to spot policy confusion.",
    "Auth & audit logs for the staff-facing version.",
])
footer(s, 12)

# ============================================================== SLIDE 13
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
tf = textbox(s, Inches(1.0), Inches(0.7), Inches(11.3), Inches(1.0))
set_run(tf.paragraphs[0].add_run(), "Q & A", 40, WHITE, bold=True)
qa = [
    ('"Can it make mistakes?"',
     "It only answers from official docs and cites them; if unsure, it says so."),
    ('"What does it cost to run?"',
     "Pay-per-use AI API; small per-query cost, no infrastructure to maintain."),
    ('"How do we add the latest circular?"',
     "One command, same day."),
    ('"Is our data safe?"',
     "Documents are public circulars; no consumer personal data is stored."),
]
tf2 = textbox(s, Inches(1.0), Inches(1.9), Inches(11.3), Inches(4.0))
for i, (q, a) in enumerate(qa):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.space_after = Pt(4)
    set_run(p.add_run(), q, 19, RGBColor(0x9F, 0xC2, 0xE0), bold=True)
    pa = tf2.add_paragraph()
    pa.space_after = Pt(14)
    set_run(pa.add_run(), a, 17, WHITE)
tf3 = textbox(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(1.0))
set_run(tf3.paragraphs[0].add_run(), "Thank you  —  Questions?", 24,
        RGBColor(0x8E, 0xE6, 0xB0), bold=True)

prs.save("UHBVN_Assistant_Demo.pptx")
print("Saved UHBVN_Assistant_Demo.pptx with", len(prs.slides._sldIdLst), "slides")
